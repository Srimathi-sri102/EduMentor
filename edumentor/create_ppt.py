from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(99, 102, 241)
SECONDARY = RGBColor(139, 92, 246)
CYAN = RGBColor(14, 165, 233)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
DARK = RGBColor(30, 41, 59)
TEXT = RGBColor(51, 65, 85)
TEXT2 = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
SURFACE = RGBColor(248, 250, 252)
LIGHT_PRIMARY = RGBColor(238, 238, 255)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, bold=False, color=TEXT, space_before=6, space_after=2, bullet=False, font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet:
        p.level = 0
    return p

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, num, total=20):
    add_text_box(slide, 11.5, 7.0, 1.5, 0.4, f"{num} / {total}", 11, False, TEXT2, PP_ALIGN.RIGHT)

def add_footer(slide):
    add_text_box(slide, 0.5, 7.0, 2, 0.4, "EduMentor", 11, True, PRIMARY)

# =================== SLIDE 1: TITLE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 7.5, LIGHT_PRIMARY)
add_text_box(slide, 1, 1.2, 11.3, 1, "🎓", 60, False, DARK, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.3, 11.3, 1, "EduMentor", 54, True, PRIMARY, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.5, 10.3, 0.8, "AI-Powered Course Generation & Personalized Learning Platform", 24, False, TEXT, PP_ALIGN.CENTER)
add_text_box(slide, 2, 4.6, 9.3, 0.5, "Flask  •  Groq LLM  •  Llama 3.3 70B  •  SQLite  •  Piston API", 16, False, TEXT2, PP_ALIGN.CENTER)
add_text_box(slide, 2, 5.8, 9.3, 0.5, "Project Presentation  •  2026", 14, False, TEXT2, PP_ALIGN.CENTER)
add_slide_number(slide, 1)

# =================== SLIDE 2: PROBLEM STATEMENT ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, PRIMARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "📌 Problem Statement", 34, True, DARK)
add_text_box(slide, 0.7, 1.4, 11, 0.5, "Traditional learning platforms suffer from critical challenges:", 18, True, TEXT)

cards = [
    ("📦 Static, One-Size-Fits-All Courses", "Fixed curricula that cannot adapt to individual learner's pace, prior knowledge, or specific goals. Beginners and advanced learners receive identical content structure.", 0.7, 2.1),
    ("🚫 No Dynamic Course Creation", "Creating quality courses requires significant manual effort by educators. Niche topics have very limited options available on existing platforms.", 6.5, 2.1),
    ("🔀 Fragmented Learning Tools", "Students juggle multiple disconnected platforms — one for courses, another for coding practice, another for quizzes, and another for interview prep.", 0.7, 4.3),
    ("🎯 Inadequate Interview Preparation", "Aspiring developers lack structured, realistic interview practice with immediate AI feedback, often relying on expensive coaching or scattered free resources.", 6.5, 4.3),
]
for title, desc, l, t in cards:
    add_rect(slide, l, t, 5.5, 1.9, SURFACE, RGBColor(226, 232, 240))
    add_text_box(slide, l+0.3, t+0.2, 4.9, 0.5, title, 16, True, DARK)
    add_text_box(slide, l+0.3, t+0.8, 4.9, 1.0, desc, 13, False, TEXT2)
add_footer(slide); add_slide_number(slide, 2)

# =================== SLIDE 3: PROPOSED SOLUTION ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, SECONDARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "💡 Proposed Solution – EduMentor", 34, True, DARK)
add_rect(slide, 0.7, 1.5, 11.9, 0.9, RGBColor(243, 243, 255), PRIMARY)
tf = add_text_box(slide, 1.0, 1.55, 11.3, 0.8, "EduMentor is an AI-powered web platform built around Dynamic Course Generation using LLMs.", 16, False, TEXT)
add_para(tf, "It generates complete, structured courses tailored to any skill, level, and duration — in seconds.", 16, False, TEXT, 2)

modules = [
    ("📚", "AI Course\nGeneration", "Personalized module-by-module\ncourses with dynamic lessons", 0.7, PRIMARY),
    ("💻", "Smart Coding\nIDE", "AI-generated problems with\nmulti-engine code execution", 2.8, SECONDARY),
    ("📝", "Adaptive\nQuizzes", "MCQ tests with instant AI\ngrading and explanations", 4.9, CYAN),
    ("🎯", "Interview\nPrep", "Mock interviews with AI\nevaluation and scoring", 7.0, GREEN),
    ("🗺️", "Learning\nRoadmaps", "Week-by-week structured\nlearning plans", 9.1, AMBER),
    ("📈", "Progress\nDashboard", "Unified analytics across\nall modules", 11.2, RED),
]
for icon, title, desc, l, c in modules:
    add_rect(slide, l, 2.8, 1.85, 3.2, SURFACE, c)
    add_text_box(slide, l, 3.0, 1.85, 0.5, icon, 30, False, DARK, PP_ALIGN.CENTER)
    add_text_box(slide, l+0.1, 3.6, 1.65, 0.7, title, 13, True, DARK, PP_ALIGN.CENTER)
    add_text_box(slide, l+0.1, 4.5, 1.65, 1.0, desc, 10, False, TEXT2, PP_ALIGN.CENTER)
add_footer(slide); add_slide_number(slide, 3)

# =================== SLIDE 4: OBJECTIVES ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, CYAN)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🎯 Project Objectives", 34, True, DARK)

objectives = [
    "Build an AI-powered course generation system that creates personalized, structured courses for any technical skill in seconds",
    "Provide an integrated environment where learners can study, code, test, and prepare for interviews — from a single platform",
    "Leverage free and open-source AI services (Groq LLM with Llama 3.3 70B, Piston API) for cost-effective accessibility",
    "Implement dynamic lesson content generation — lessons are only generated when the user clicks on them, reducing upfront compute",
    "Track and visualize learner progress with aggregated analytics to encourage consistency and self-improvement",
    "Demonstrate the practical application of Large Language Models in the Education Technology domain",
    "Support multi-engine code execution with cascading fallback strategy (Piston → Judge0 → AI Simulation)",
]
tf = add_text_box(slide, 0.9, 1.5, 11.5, 0.4, "", 14, False, TEXT)
for obj in objectives:
    add_para(tf, f"✅  {obj}", 15, False, TEXT, 8, 4)
add_rect(slide, 0.7, 6.0, 11.9, 0.7, RGBColor(243, 243, 255), PRIMARY)
add_text_box(slide, 1.0, 6.05, 11.3, 0.6, "Key Innovation: Every course, problem, quiz, and interview question is generated dynamically by AI — no pre-authored static content needed.", 14, True, PRIMARY)
add_footer(slide); add_slide_number(slide, 4)

# =================== SLIDE 5: TECH STACK ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, PRIMARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "⚙️ Technology Stack Overview", 34, True, DARK)

# Table header
add_rect(slide, 0.7, 1.5, 11.9, 0.5, PRIMARY)
for i, h in enumerate(["Layer", "Technology", "Version", "Purpose"]):
    ws = [2.2, 3.5, 1.5, 4.7]
    ls = [0.7, 2.9, 6.4, 7.9]
    add_text_box(slide, ls[i], 1.5, ws[i], 0.5, h, 14, True, WHITE)

rows = [
    ("Backend", "Python + Flask", "3.x / 3.0.0", "Micro web framework with Blueprints"),
    ("Database", "SQLite + SQLAlchemy", "3.1.1", "ORM-based data persistence"),
    ("Authentication", "Flask-Login + Flask-Bcrypt", "0.6.3 / 1.0.1", "Sessions, bcrypt password hashing"),
    ("AI / LLM", "Groq API (Llama 3.3 70B)", "0.11.0", "All AI content generation"),
    ("Code Execution", "Piston API + Judge0", "—", "Sandboxed execution (8 languages)"),
    ("Frontend", "HTML5, CSS3, Vanilla JS", "—", "Jinja2 templating, custom design system"),
    ("Typography", "Google Fonts", "—", "Inter (UI) + JetBrains Mono (code)"),
    ("Environment", "python-dotenv", "1.0.1", "Secure API key management via .env"),
]
for ri, (a,b,c,d) in enumerate(rows):
    y = 2.05 + ri*0.55
    bg = SURFACE if ri%2==0 else WHITE
    add_rect(slide, 0.7, y, 11.9, 0.5, bg)
    add_text_box(slide, 0.8, y, 2.0, 0.5, a, 13, True, DARK)
    add_text_box(slide, 2.9, y, 3.5, 0.5, b, 13, False, TEXT)
    add_text_box(slide, 6.4, y, 1.5, 0.5, c, 13, False, TEXT2)
    add_text_box(slide, 7.9, y, 4.7, 0.5, d, 13, False, TEXT2)
add_footer(slide); add_slide_number(slide, 5)

# =================== SLIDE 6: SYSTEM ARCHITECTURE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, SECONDARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🏗️ System Architecture", 34, True, DARK)

add_rect(slide, 4.5, 1.4, 4.3, 0.7, SURFACE, PRIMARY)
add_text_box(slide, 4.5, 1.45, 4.3, 0.6, "👤 User (Browser)", 18, True, PRIMARY, PP_ALIGN.CENTER)
add_text_box(slide, 6.2, 2.1, 1, 0.4, "⬇ ⬆", 16, True, PRIMARY, PP_ALIGN.CENTER)
add_rect(slide, 3, 2.5, 7.3, 0.7, SURFACE, SECONDARY)
add_text_box(slide, 3, 2.55, 7.3, 0.6, "Flask Application Server (Werkzeug WSGI)", 18, True, SECONDARY, PP_ALIGN.CENTER)

arch_items = [
    ("📄 Jinja2 Templates", "9 HTML pages", 0.7, CYAN),
    ("🔗 Flask Blueprints", "8 route modules", 4.5, PRIMARY),
    ("🗃️ SQLAlchemy ORM", "7 database models", 8.5, GREEN),
]
for title, sub, l, c in arch_items:
    add_rect(slide, l, 3.7, 3.5, 1.0, SURFACE, c)
    add_text_box(slide, l, 3.8, 3.5, 0.5, title, 15, True, DARK, PP_ALIGN.CENTER)
    add_text_box(slide, l, 4.25, 3.5, 0.4, sub, 12, False, TEXT2, PP_ALIGN.CENTER)

arch_items2 = [
    ("🤖 Groq LLM API", "Llama 3.3 70B Versatile", 0.7, AMBER),
    ("⚡ Piston / Judge0", "Code Execution Engines", 4.5, RGBColor(236,72,153)),
    ("💾 SQLite Database", "edumentor.db", 8.5, GREEN),
]
for title, sub, l, c in arch_items2:
    add_rect(slide, l, 5.2, 3.5, 1.0, SURFACE, c)
    add_text_box(slide, l, 5.3, 3.5, 0.5, title, 15, True, DARK, PP_ALIGN.CENTER)
    add_text_box(slide, l, 5.75, 3.5, 0.4, sub, 12, False, TEXT2, PP_ALIGN.CENTER)

add_text_box(slide, 1, 6.5, 11.3, 0.4, "Pattern: Application Factory (create_app) → Modular Blueprints → RESTful JSON APIs → Async Frontend Fetch Calls", 13, True, TEXT2, PP_ALIGN.CENTER)
add_footer(slide); add_slide_number(slide, 6)

# =================== SLIDE 7: PROJECT STRUCTURE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, CYAN)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "📁 Project Directory Structure", 34, True, DARK)

add_rect(slide, 0.7, 1.5, 5.5, 5.0, RGBColor(30,41,59))
structure = """edumentor/
├── app.py             # App factory (create_app)
├── config.py          # Config class (.env loading)
├── models.py          # 7 database models
├── requirements.txt   # 8 Python dependencies
├── .env               # API keys (secret)
├── routes/
│   ├── auth.py        # Login / Register / Logout
│   ├── main.py        # Dashboard routing
│   ├── roadmap.py     # Roadmap CRUD + AI
│   ├── course.py      # Course generation
│   ├── ide.py         # Code execution + eval
│   ├── quiz.py        # Quiz engine
│   ├── interview.py   # Interview prep
│   └── progress.py    # Analytics API
├── templates/         # 9 Jinja2 HTML files
├── static/css/        # Custom design system
├── static/js/         # Client-side utilities
└── instance/          # SQLite database"""
add_text_box(slide, 0.9, 1.6, 5.1, 4.8, structure, 11, False, RGBColor(226,232,240), font_name='Consolas')

descs = [
    ("🏭 Application Factory Pattern", "create_app() in app.py initializes Flask, registers all 8 Blueprints, configures extensions (SQLAlchemy, Login, Bcrypt), and creates database tables automatically."),
    ("🔌 Blueprint Architecture", "Each feature is a self-contained Flask Blueprint with its own routes and API endpoints. Total: 8 blueprints with 19 routes across all modules."),
    ("🎨 Template Inheritance", "base.html provides the master layout with sidebar navigation and topbar. All 8 pages extend it using Jinja2's {% extends 'base.html' %} pattern."),
]
for i, (t, d) in enumerate(descs):
    y = 1.5 + i*1.65
    add_rect(slide, 6.6, y, 5.8, 1.5, SURFACE, RGBColor(226,232,240))
    add_text_box(slide, 6.8, y+0.15, 5.4, 0.4, t, 15, True, DARK)
    add_text_box(slide, 6.8, y+0.6, 5.4, 0.8, d, 12, False, TEXT2)
add_footer(slide); add_slide_number(slide, 7)

# =================== SLIDE 8: DATABASE MODELS ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, PRIMARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🗄️ Database Models (SQLAlchemy ORM)", 34, True, DARK)

headers = ["Model", "Table", "Key Columns", "Relationships"]
ws = [1.8, 2.0, 4.8, 3.3]
ls = [0.7, 2.5, 4.5, 9.3]
add_rect(slide, 0.7, 1.4, 11.9, 0.5, PRIMARY)
for i, h in enumerate(headers):
    add_text_box(slide, ls[i], 1.4, ws[i], 0.5, h, 13, True, WHITE)

models = [
    ("User", "users", "id, username, email, password_hash, skill, level", "→ roadmaps, sessions, quizzes, courses"),
    ("Roadmap", "roadmaps", "id, user_id, skill, level, content (JSON)", "← User (FK)"),
    ("CodingSession", "coding_sessions", "id, user_id, problem_title, user_code, result", "← User (FK)"),
    ("Quiz", "quizzes", "id, user_id, skill, questions (JSON), score", "← User (FK)"),
    ("InterviewSession", "interview_sessions", "id, user_id, skill, level, questions (JSON)", "← User (FK)"),
    ("Course", "courses", "id, user_id, title, structure (JSON), completed", "← User, → LessonContent"),
    ("LessonContent", "lesson_contents", "id, course_id, module_id, lesson_id, content", "← Course (FK, cascade)"),
]
for ri, (a,b,c,d) in enumerate(models):
    y = 1.95 + ri*0.6
    bg = SURFACE if ri%2==0 else WHITE
    add_rect(slide, 0.7, y, 11.9, 0.55, bg)
    add_text_box(slide, 0.8, y, 1.6, 0.5, a, 12, True, DARK)
    add_text_box(slide, 2.5, y, 2.0, 0.5, b, 12, False, TEXT)
    add_text_box(slide, 4.5, y, 4.8, 0.5, c, 11, False, TEXT2)
    add_text_box(slide, 9.3, y, 3.3, 0.5, d, 11, False, TEXT2)

add_rect(slide, 0.7, 6.3, 11.9, 0.7, RGBColor(243,243,255), PRIMARY)
add_text_box(slide, 1.0, 6.35, 11.3, 0.6, "Design: Complex AI-generated data stored as serialized JSON in TEXT columns. Flexible schema for dynamic AI output while maintaining relational integrity via foreign keys.", 13, False, PRIMARY)
add_footer(slide); add_slide_number(slide, 8)

# =================== SLIDE 9: AI/LLM INTEGRATION ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, SECONDARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🤖 AI / LLM Integration (Groq + Llama 3.3)", 34, True, DARK)

add_rect(slide, 0.7, 1.5, 5.8, 2.0, SURFACE, RGBColor(226,232,240))
add_text_box(slide, 0.9, 1.6, 5.4, 0.4, "⚡ Why Groq?", 16, True, DARK)
add_text_box(slide, 0.9, 2.1, 5.4, 1.2, "Groq's custom LPU™ hardware delivers the fastest LLM inference available. The free tier with Llama 3.3 70B (70 billion parameters) provides high-quality, cost-effective AI generation for all platform features.", 13, False, TEXT2)

add_rect(slide, 0.7, 3.7, 5.8, 2.0, SURFACE, RGBColor(226,232,240))
add_text_box(slide, 0.9, 3.8, 5.4, 0.4, "📋 Structured JSON Output", 16, True, DARK)
add_text_box(slide, 0.9, 4.3, 5.4, 1.2, "All AI prompts request structured JSON responses. A _clean_json() utility strips markdown code fences before json.loads() parsing. System prompts enforce 'Return valid JSON only, no markdown.'", 13, False, TEXT2)

# AI features table
add_rect(slide, 6.8, 1.5, 5.8, 0.5, SECONDARY)
add_text_box(slide, 6.9, 1.5, 2.7, 0.5, "Feature", 13, True, WHITE)
add_text_box(slide, 9.6, 1.5, 3.0, 0.5, "AI Role", 13, True, WHITE)
ai_rows = [
    ("Course Generation", "Creates full curriculum outlines"),
    ("Lesson Generation", "Dynamic lesson content on-demand"),
    ("Roadmap Generation", "Week-by-week learning plans"),
    ("Problem Generation", "Coding challenges with starter code"),
    ("Code Evaluation", "Scores, feedback, complexity analysis"),
    ("Quiz Generation", "MCQs with answer explanations"),
    ("Interview Questions", "Technical / behavioral Q&A"),
    ("Answer Evaluation", "Grades, strengths, improvements"),
]
for ri, (a,b) in enumerate(ai_rows):
    y = 2.05 + ri*0.5
    bg = SURFACE if ri%2==0 else WHITE
    add_rect(slide, 6.8, y, 5.8, 0.45, bg)
    add_text_box(slide, 6.9, y, 2.7, 0.45, a, 12, True if ri==0 else False, DARK)
    add_text_box(slide, 9.6, y, 3.0, 0.45, b, 12, False, TEXT2)
add_footer(slide); add_slide_number(slide, 9)

# =================== SLIDE 10: CORE FEATURE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, CYAN)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "📚 Core Feature: Dynamic Course Generator", 34, True, DARK)
add_text_box(slide, 0.7, 1.35, 11, 0.4, "Two-phase generation strategy to minimize compute and maximize responsiveness:", 17, True, TEXT)

add_rect(slide, 0.7, 1.9, 5.8, 3.0, SURFACE, PRIMARY)
add_text_box(slide, 0.9, 2.0, 5.4, 0.4, "Phase 1 — Outline Generation", 16, True, PRIMARY)
tf = add_text_box(slide, 0.9, 2.5, 5.4, 0.4, "• User inputs: skill, level, module count (3–10)", 13, False, TEXT)
add_para(tf, "• AI generates full curriculum structure as JSON", 13, False, TEXT, 4)
add_para(tf, "• Each lesson: title, type (theory/practice/project), duration", 13, False, TEXT, 4)
add_para(tf, "• Rendered as interactive tree with collapsible modules", 13, False, TEXT, 4)
add_para(tf, "• Endpoint: POST /api/course/generate", 13, True, SECONDARY, 4)

add_rect(slide, 6.8, 1.9, 5.8, 3.0, SURFACE, GREEN)
add_text_box(slide, 7.0, 2.0, 5.4, 0.4, "Phase 2 — Lazy Lesson Generation", 16, True, GREEN)
tf = add_text_box(slide, 7.0, 2.5, 5.4, 0.4, "• Lesson content NOT generated upfront", 13, False, TEXT)
add_para(tf, "• Generated on-demand when user clicks a lesson", 13, False, TEXT, 4)
add_para(tf, "• Content: text, code examples, exercises, takeaways", 13, False, TEXT, 4)
add_para(tf, "• Saved courses cache content in LessonContent table", 13, False, TEXT, 4)
add_para(tf, "• Endpoint: POST /api/course/generate-lesson", 13, True, GREEN, 4)

# Flow
flow_items = ["📝 User Input", "→", "🤖 AI Outline", "→", "🖱️ Click Lesson", "→", "⚡ AI Content", "→", "💾 Cache"]
x = 1.0
for item in flow_items:
    if item == "→":
        add_text_box(slide, x, 5.3, 0.5, 0.5, "→", 22, True, PRIMARY, PP_ALIGN.CENTER)
        x += 0.5
    else:
        add_rect(slide, x, 5.2, 1.8, 0.6, SURFACE, RGBColor(226,232,240))
        add_text_box(slide, x, 5.25, 1.8, 0.5, item, 12, True, DARK, PP_ALIGN.CENTER)
        x += 2.1
add_footer(slide); add_slide_number(slide, 10)

# =================== SLIDE 11: CODING IDE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, PRIMARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "💻 Smart Coding IDE Module", 34, True, DARK)

engines = [
    ("🥇 Priority 1: Piston API", "Free, no API key needed. Sandboxed execution via emkc.org. Supports 8 languages with compile + run timeouts.", GREEN),
    ("🥈 Priority 2: Judge0 (RapidAPI)", "Used when Piston fails. Base64-encoded I/O. Requires JUDGE0_API_KEY configured in .env file.", AMBER),
    ("🥉 Priority 3: AI Simulation", "Groq LLM 'mentally executes' code and returns expected output. Ultimate last-resort fallback.", RGBColor(236,72,153)),
]
for i, (t, d, c) in enumerate(engines):
    y = 1.5 + i*1.3
    add_rect(slide, 0.7, y, 6.0, 1.1, SURFACE, c)
    add_text_box(slide, 0.9, y+0.1, 5.6, 0.4, t, 15, True, DARK)
    add_text_box(slide, 0.9, y+0.5, 5.6, 0.5, d, 12, False, TEXT2)

add_text_box(slide, 7.0, 1.5, 5.5, 0.4, "Supported Languages (8)", 16, True, DARK)
langs = ["🐍 Python", "📜 JavaScript", "☕ Java", "⚙️ C++", "🔧 C", "🦫 Go", "🦀 Rust", "📘 TypeScript"]
for i, lang in enumerate(langs):
    row, col = divmod(i, 2)
    add_rect(slide, 7.0+col*2.7, 2.0+row*0.5, 2.5, 0.4, SURFACE, RGBColor(226,232,240))
    add_text_box(slide, 7.1+col*2.7, 2.0+row*0.5, 2.3, 0.4, lang, 12, False, DARK, PP_ALIGN.CENTER)

add_text_box(slide, 7.0, 4.2, 5.5, 0.4, "AI Evaluation JSON Output", 15, True, DARK)
add_rect(slide, 7.0, 4.6, 5.5, 1.8, RGBColor(30,41,59))
eval_json = '{\n  "score": 85,\n  "correct": true,\n  "feedback": "Good solution...",\n  "improvements": ["Use const..."],\n  "time_complexity": "O(n)",\n  "space_complexity": "O(1)"\n}'
add_text_box(slide, 7.2, 4.7, 5.1, 1.6, eval_json, 11, False, RGBColor(226,232,240), font_name='Consolas')
add_footer(slide); add_slide_number(slide, 11)

# =================== SLIDE 12: QUIZ ENGINE ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, SECONDARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "📝 Adaptive Quiz Engine", 34, True, DARK)

tf = add_text_box(slide, 0.7, 1.5, 6.0, 0.4, "How It Works", 18, True, DARK)
add_para(tf, "• User selects skill, level, optional topic, question count (up to 20)", 14, False, TEXT, 8)
add_para(tf, "• AI generates MCQs with 4 options (A–D), correct answer, explanation", 14, False, TEXT, 6)
add_para(tf, "• Quiz auto-saved to database on generation for history tracking", 14, False, TEXT, 6)
add_para(tf, "• User answers submitted via POST /api/quiz/submit", 14, False, TEXT, 6)
add_para(tf, "• Server-side grading: compares user answer with correct answer", 14, False, TEXT, 6)
add_para(tf, "• Returns: score, total, percentage, per-question results", 14, False, TEXT, 6)

add_rect(slide, 0.7, 5.3, 6.0, 0.8, RGBColor(243,243,255), SECONDARY)
add_text_box(slide, 0.9, 5.35, 5.6, 0.7, "Data Flow: Generate → Save to DB → User answers → Submit → Server grades → Update score → Return results with explanations", 12, True, SECONDARY)

add_text_box(slide, 7.0, 1.5, 5.5, 0.4, "AI Response Structure", 16, True, DARK)
add_rect(slide, 7.0, 2.0, 5.5, 2.5, RGBColor(30,41,59))
quiz_json = '[{\n  "id": 1,\n  "question": "What is a closure?",\n  "options": [\n    "A) A function...",\n    "B) A variable...",\n    "C) A class...",\n    "D) An object..."\n  ],\n  "correct_answer": "A",\n  "explanation": "A closure is..."\n}]'
add_text_box(slide, 7.2, 2.1, 5.1, 2.3, quiz_json, 11, False, RGBColor(226,232,240), font_name='Consolas')

# API table
add_text_box(slide, 7.0, 4.8, 5.5, 0.4, "API Endpoints", 15, True, DARK)
add_rect(slide, 7.0, 5.2, 5.5, 0.35, SECONDARY)
add_text_box(slide, 7.1, 5.2, 1.5, 0.35, "Method", 12, True, WHITE)
add_text_box(slide, 8.6, 5.2, 3.8, 0.35, "Endpoint", 12, True, WHITE)
eps = [("GET", "/quiz"), ("POST", "/api/quiz/generate"), ("POST", "/api/quiz/submit")]
for i, (m, e) in enumerate(eps):
    y = 5.55 + i*0.35
    add_rect(slide, 7.0, y, 5.5, 0.35, SURFACE if i%2==0 else WHITE)
    add_text_box(slide, 7.1, y, 1.5, 0.35, m, 12, True, DARK)
    add_text_box(slide, 8.6, y, 3.8, 0.35, e, 12, False, TEXT)
add_footer(slide); add_slide_number(slide, 12)

# =================== SLIDE 13: INTERVIEW PREP ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, CYAN)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🎯 AI Interview Preparation", 34, True, DARK)

tf = add_text_box(slide, 0.7, 1.5, 6.0, 0.4, "Features", 18, True, DARK)
add_para(tf, "• Generates Technical, Behavioral, and System Design questions", 14, False, TEXT, 8)
add_para(tf, "• Up to 15 questions per session, fully configurable", 14, False, TEXT, 6)
add_para(tf, "• Each question: model answer, key points, follow-up questions", 14, False, TEXT, 6)
add_para(tf, "• Users submit their own answers for AI evaluation", 14, False, TEXT, 6)
add_para(tf, "• AI scores on a 0–10 scale with grade, strengths, improvements", 14, False, TEXT, 6)

add_rect(slide, 7.0, 1.5, 5.5, 0.7, SURFACE, CYAN)
add_text_box(slide, 7.2, 1.55, 5.1, 0.6, "1. Generate → POST /api/interview/generate\n   Groq LLM creates questions → Saved to DB", 12, False, TEXT)
add_rect(slide, 7.0, 2.4, 5.5, 0.7, SURFACE, CYAN)
add_text_box(slide, 7.2, 2.45, 5.1, 0.6, "2. Practice → User reads questions, views\n   model answers, prepares responses", 12, False, TEXT)
add_rect(slide, 7.0, 3.3, 5.5, 0.7, SURFACE, CYAN)
add_text_box(slide, 7.2, 3.35, 5.1, 0.6, "3. Evaluate → POST /api/interview/evaluate-answer\n   User answer + model answer → LLM → Score", 12, False, TEXT)

add_text_box(slide, 0.7, 4.5, 5.5, 0.4, "Evaluation JSON Output", 15, True, DARK)
add_rect(slide, 0.7, 4.9, 5.5, 2.0, RGBColor(30,41,59))
eval_json2 = '{\n  "score": 8,\n  "grade": "Good",\n  "feedback": "Strong understanding...",\n  "strengths": ["Clear explanation"],\n  "improvements": ["Add examples"]\n}'
add_text_box(slide, 0.9, 5.0, 5.1, 1.8, eval_json2, 12, False, RGBColor(226,232,240), font_name='Consolas')
add_footer(slide); add_slide_number(slide, 13)

# =================== SLIDE 14: ROADMAP + PROGRESS ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, GREEN)
add_text_box(slide, 0.7, 0.6, 12, 0.7, "🗺️ Learning Roadmaps & 📈 Progress Dashboard", 32, True, DARK)

add_rect(slide, 0.7, 1.5, 5.8, 2.8, SURFACE, RGBColor(226,232,240))
add_text_box(slide, 0.9, 1.6, 5.4, 0.4, "Roadmap Generator", 16, True, DARK)
tf = add_text_box(slide, 0.9, 2.1, 5.4, 0.4, "• User inputs: skill, level, duration (weeks)", 13, False, TEXT)
add_para(tf, "• AI generates week-by-week learning plan", 13, False, TEXT, 4)
add_para(tf, "• Each week: topics, resources, projects, goals", 13, False, TEXT, 4)
add_para(tf, "• Full CRUD: generate, save, load, delete", 13, False, TEXT, 4)
add_para(tf, "• Stored as serialized JSON in Roadmap table", 13, False, TEXT, 4)

add_rect(slide, 6.8, 1.5, 5.8, 2.8, SURFACE, RGBColor(226,232,240))
add_text_box(slide, 7.0, 1.6, 5.4, 0.4, "Progress Dashboard", 16, True, DARK)
tf = add_text_box(slide, 7.0, 2.1, 5.4, 0.4, "• Aggregated stats across all modules", 13, False, TEXT)
add_para(tf, "• Weekly quiz score trends (last 8 weeks)", 13, False, TEXT, 4)
add_para(tf, "• Activity breakdown by module type", 13, False, TEXT, 4)
add_para(tf, "• Coding sessions pass rate tracking", 13, False, TEXT, 4)
add_para(tf, "• Average quiz score calculation", 13, False, TEXT, 4)

metrics = [
    ("Roadmaps", "Count created", PRIMARY),
    ("Coding", "Sessions + pass rate", SECONDARY),
    ("Quizzes", "Avg score + trends", CYAN),
    ("Interviews", "Session count", GREEN),
]
for i, (t, d, c) in enumerate(metrics):
    x = 0.7 + i*3.1
    add_rect(slide, x, 4.8, 2.8, 1.0, SURFACE, c)
    add_text_box(slide, x, 4.9, 2.8, 0.4, t, 15, True, c, PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.3, 2.8, 0.4, d, 12, False, TEXT2, PP_ALIGN.CENTER)
add_footer(slide); add_slide_number(slide, 14)

# =================== SLIDE 15: AUTHENTICATION ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, SECONDARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🔐 Authentication & Security", 34, True, DARK)

auth_cards = [
    ("Registration", "Username + email uniqueness check → Password hashed with bcrypt + salt → User created → Auto-login via Flask-Login session cookie"),
    ("Login", "Username lookup → bcrypt.check_password_hash() verification → Session created with 'Remember Me' → Redirect to dashboard"),
    ("Route Protection", "All protected routes decorated with @login_required. Unauthorized access redirects to login page with flash message."),
]
for i, (t, d) in enumerate(auth_cards):
    y = 1.4 + i*1.2
    add_rect(slide, 0.7, y, 6.0, 1.0, SURFACE, RGBColor(226,232,240))
    add_text_box(slide, 0.9, y+0.1, 5.6, 0.3, t, 15, True, DARK)
    add_text_box(slide, 0.9, y+0.45, 5.6, 0.5, d, 12, False, TEXT2)

add_text_box(slide, 7.0, 1.4, 5.5, 0.4, "Security Measures", 17, True, DARK)
sec = [
    "✅  Passwords hashed with bcrypt (salt + iterative hashing)",
    "✅  API keys stored in .env file, never in source code",
    "✅  Flask SECRET_KEY for session cookie encryption",
    "✅  All API endpoints require authenticated session",
    "✅  User data isolation via user_id foreign keys",
    "✅  CSRF protection via Flask session tokens",
]
tf = add_text_box(slide, 7.0, 1.9, 5.5, 0.4, sec[0], 14, False, TEXT)
for s in sec[1:]:
    add_para(tf, s, 14, False, TEXT, 8)
add_footer(slide); add_slide_number(slide, 15)

# =================== SLIDE 16: FRONTEND DESIGN ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, CYAN)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🎨 Frontend Design System", 34, True, DARK)

tf = add_text_box(slide, 0.7, 1.4, 6.0, 0.4, "Custom CSS Design System", 17, True, DARK)
add_para(tf, "• CSS Custom Properties (variables) for consistent theming", 13, False, TEXT, 8)
add_para(tf, "• Dark theme with glassmorphism (backdrop-filter: blur)", 13, False, TEXT, 5)
add_para(tf, "• Animated floating orbs on landing page", 13, False, TEXT, 5)
add_para(tf, "• Responsive grid layouts with @media breakpoints", 13, False, TEXT, 5)
add_para(tf, "• Component library: cards, buttons, forms, badges, toasts", 13, False, TEXT, 5)
add_para(tf, "• Smooth animations: fadeInUp, slideDown, spin, float", 13, False, TEXT, 5)

add_text_box(slide, 0.7, 4.3, 6.0, 0.4, "Color Palette", 15, True, DARK)
colors = [
    ("#6366f1", "Primary", PRIMARY), ("#8b5cf6", "Secondary", SECONDARY),
    ("#0ea5e9", "Cyan", CYAN), ("#10b981", "Green", GREEN),
    ("#f59e0b", "Amber", AMBER), ("#ef4444", "Red", RED),
    ("#ec4899", "Pink", RGBColor(236,72,153)), ("#050510", "Background", RGBColor(5,5,16)),
]
for i, (hex_c, name, c) in enumerate(colors):
    x = 0.7 + i*0.95
    add_rect(slide, x, 4.8, 0.7, 0.7, c)
    add_text_box(slide, x-0.1, 5.55, 0.9, 0.3, name, 9, False, TEXT2, PP_ALIGN.CENTER)

# Template list
add_text_box(slide, 7.0, 1.4, 5.5, 0.4, "Template Architecture (Jinja2)", 16, True, DARK)
templates = ["base.html → Master layout (sidebar + topbar)", "index.html → Landing page / Auth forms",
    "dashboard.html → User dashboard with stats", "roadmap.html → Roadmap generator UI",
    "course.html → Dynamic course generator", "ide.html → Coding IDE interface",
    "quiz.html → Quiz engine UI", "interview.html → Interview prep UI", "progress.html → Analytics dashboard"]
tf = add_text_box(slide, 7.0, 1.9, 5.5, 0.3, templates[0], 12, True, TEXT)
for t in templates[1:]:
    add_para(tf, t, 12, False, TEXT, 4)

add_text_box(slide, 7.0, 4.8, 5.5, 0.4, "JavaScript (Vanilla — No Frameworks)", 15, True, DARK)
tf = add_text_box(slide, 7.0, 5.2, 5.5, 0.3, "• Toast notification system", 12, False, TEXT)
add_para(tf, "• AJAX via fetch() — zero dependencies", 12, False, TEXT, 4)
add_para(tf, "• Sidebar toggle for mobile responsiveness", 12, False, TEXT, 4)
add_para(tf, "• Tab key insertion in code editor", 12, False, TEXT, 4)
add_footer(slide); add_slide_number(slide, 16)

# =================== SLIDE 17: API ENDPOINTS ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, PRIMARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🔗 API Endpoints Summary", 34, True, DARK)

add_rect(slide, 0.7, 1.4, 11.9, 0.5, PRIMARY)
for i, h in enumerate(["Module", "Method", "Endpoint", "Description"]):
    ws2 = [1.8, 1.2, 4.2, 4.7]
    ls2 = [0.7, 2.5, 3.7, 7.9]
    add_text_box(slide, ls2[i], 1.4, ws2[i], 0.5, h, 13, True, WHITE)

api_rows = [
    ("Auth", "POST", "/login", "User login with session"),
    ("", "POST", "/register", "Create account + auto-login"),
    ("", "GET", "/logout", "Destroy session"),
    ("Course", "POST", "/api/course/generate", "AI generates curriculum outline"),
    ("", "POST", "/api/course/generate-lesson", "Dynamic lesson content generation"),
    ("", "POST", "/api/course/save", "Save course to database"),
    ("IDE", "POST", "/api/ide/generate-problem", "AI generates coding problem"),
    ("", "POST", "/api/ide/execute", "Run code (Piston→Judge0→AI)"),
    ("", "POST", "/api/ide/evaluate", "AI evaluates solution quality"),
    ("Quiz", "POST", "/api/quiz/generate", "Generate MCQ questions"),
    ("", "POST", "/api/quiz/submit", "Submit answers, get scores"),
    ("Interview", "POST", "/api/interview/generate", "Generate interview questions"),
    ("Progress", "GET", "/api/progress/stats", "Get aggregated analytics"),
]
for ri, (a,b,c,d) in enumerate(api_rows):
    y = 1.95 + ri*0.4
    bg = SURFACE if ri%2==0 else WHITE
    add_rect(slide, 0.7, y, 11.9, 0.38, bg)
    add_text_box(slide, 0.8, y, 1.6, 0.38, a, 11, True, DARK)
    add_text_box(slide, 2.5, y, 1.2, 0.38, b, 11, False, TEXT)
    add_text_box(slide, 3.7, y, 4.2, 0.38, c, 11, False, SECONDARY, font_name='Consolas')
    add_text_box(slide, 7.9, y, 4.7, 0.38, d, 11, False, TEXT2)
add_footer(slide); add_slide_number(slide, 17)

# =================== SLIDE 18: INNOVATION ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, SECONDARY)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🚀 Innovation & Uniqueness", 34, True, DARK)

innovations = [
    ("🧠 AI-Generated Courses On Demand", "Complete structured courses generated in real-time by AI — no pre-authored content needed. This is the key differentiator."),
    ("♾️ Infinite Course Variety", "Supports any skill at any difficulty level. Users can generate courses on niche topics no existing platform covers."),
    ("🔄 Lazy Lesson Generation", "Lesson content generated on-click, not upfront. Reduces compute costs and response time. Cached after first view."),
    ("🏠 Unified Learning Ecosystem", "Combines courses, coding, quizzes, interviews, and progress tracking — unlike platforms that specialize in one area."),
    ("⚡ Multi-Engine Code Execution", "Cascading fallback strategy (Piston → Judge0 → AI) ensures code always runs regardless of API availability."),
    ("💰 Cost-Effective AI Stack", "Uses free Groq API with Llama 3.3 70B — one of the most capable open-source LLMs, eliminating expensive API costs."),
]
for i, (t, d) in enumerate(innovations):
    row, col = divmod(i, 3)
    x = 0.7 + col*4.1
    y = 1.5 + row*2.5
    add_rect(slide, x, y, 3.8, 2.2, SURFACE, RGBColor(226,232,240))
    add_text_box(slide, x+0.2, y+0.2, 3.4, 0.5, t, 14, True, DARK)
    add_text_box(slide, x+0.2, y+0.8, 3.4, 1.2, d, 12, False, TEXT2)
add_footer(slide); add_slide_number(slide, 18)

# =================== SLIDE 19: FUTURE ENHANCEMENTS ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, 0.7, 0.5, 3.5, CYAN)
add_text_box(slide, 0.7, 0.6, 10, 0.7, "🔮 Future Enhancements", 34, True, DARK)

add_rect(slide, 0.7, 1.5, 5.8, 3.5, SURFACE, RGBColor(226,232,240))
add_text_box(slide, 0.9, 1.6, 5.4, 0.4, "Short-Term Goals", 17, True, DARK)
short = ["✅  Course export as PDF for offline study", "✅  Gamification: badges, streaks, leaderboards",
    "✅  Collaborative learning with study groups", "✅  Video-based explanations in courses"]
tf = add_text_box(slide, 0.9, 2.2, 5.4, 0.4, short[0], 14, False, TEXT)
for s in short[1:]:
    add_para(tf, s, 14, False, TEXT, 8)

add_rect(slide, 6.8, 1.5, 5.8, 3.5, SURFACE, RGBColor(226,232,240))
add_text_box(slide, 7.0, 1.6, 5.4, 0.4, "Long-Term Vision", 17, True, DARK)
long_goals = ["✅  GitHub integration for project-based learning", "✅  Mobile app (React Native / Flutter)",
    "✅  Multi-language UI support (Hindi, Spanish, etc.)", "✅  Production: PostgreSQL + Gunicorn + Nginx"]
tf = add_text_box(slide, 7.0, 2.2, 5.4, 0.4, long_goals[0], 14, False, TEXT)
for s in long_goals[1:]:
    add_para(tf, s, 14, False, TEXT, 8)

add_rect(slide, 0.7, 5.3, 11.9, 1.0, RGBColor(243,243,255), PRIMARY)
add_text_box(slide, 1.0, 5.4, 11.3, 0.8, "Scalability Path: SQLite → PostgreSQL  |  Werkzeug → Gunicorn/Nginx  |  Session Auth → JWT Tokens  |  Single Server → Docker + Kubernetes", 15, True, PRIMARY, PP_ALIGN.CENTER)
add_footer(slide); add_slide_number(slide, 19)

# =================== SLIDE 20: THANK YOU ===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_PRIMARY)
add_text_box(slide, 1, 1.5, 11.3, 1, "🎓", 72, False, DARK, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11.3, 1, "Thank You!", 54, True, PRIMARY, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10.3, 0.8, "EduMentor — AI-Powered Course Generation & Personalized Learning Platform", 22, False, TEXT, PP_ALIGN.CENTER)
add_text_box(slide, 2, 5.2, 9.3, 0.5, "Python + Flask  •  Groq LLM  •  Llama 3.3 70B  •  SQLite  •  Piston API", 16, False, TEXT2, PP_ALIGN.CENTER)
add_text_box(slide, 2, 6.0, 9.3, 0.5, "Questions & Discussion", 20, True, DARK, PP_ALIGN.CENTER)
add_slide_number(slide, 20)

# SAVE
output_path = r'c:\Users\acer\OneDrive\Desktop\project\edumentor\EduMentor_Presentation.pptx'
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
